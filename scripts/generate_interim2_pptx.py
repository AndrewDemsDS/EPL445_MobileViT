"""Generate Interim Presentation #2 PPTX for EPL445 MobileViT project."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

UCY_BLUE   = RGBColor(0x00, 0x3A, 0x70)
ACCENT     = RGBColor(0x00, 0x8B, 0xD0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF5, 0xF9)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x1A, 0x8A, 0x42)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)

FIGURES_DIR = os.path.join(os.path.dirname(__file__), "..", "outputs", "figures")


def add_solid_fill(shape, color):
    sp = shape.fill
    sp.solid()
    sp.fore_color.rgb = color


def set_text(tf, text, size=18, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT, italic=False):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_header_bar(slide, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.35))
    add_solid_fill(bar, UCY_BLUE)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.95), Inches(13), Pt(22))
        set_text(sub.text_frame, subtitle_text, size=14, color=ACCENT, italic=True)


def add_footer(slide, slide_num, total):
    footer = slide.shapes.add_textbox(Inches(0), Inches(7.0), Inches(13.33), Inches(0.4))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"EPL 445 – Digital Video Processing  |  University of Cyprus, Spring 2026  |  Slide {slide_num}/{total}"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    line = slide.shapes.add_shape(1, Inches(0), Inches(6.95), Inches(13.33), Pt(1.5))
    add_solid_fill(line, ACCENT)
    line.line.fill.background()


def metric_card(slide, left, top, width, height, label, value, color):
    card = slide.shapes.add_shape(1, left, top, width, height)
    add_solid_fill(card, color)
    card.line.fill.background()
    tb = slide.shapes.add_textbox(left + Pt(8), top + Pt(8), width - Pt(16), Inches(0.3))
    set_text(tb.text_frame, label, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    tb2 = slide.shapes.add_textbox(left + Pt(4), top + Inches(0.32), width - Pt(8), Inches(0.45))
    set_text(tb2.text_frame, value, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
TOTAL = 8
blank_layout = prs.slide_layouts[6]


# SLIDE 1 — Title
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, UCY_BLUE); bg.line.fill.background()
stripe = sl.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(0.08))
add_solid_fill(stripe, ACCENT); stripe.line.fill.background()
tb = sl.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.2))
set_text(tb.text_frame, "MobileViT-Based Traffic Analysis", size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(11.5), Inches(0.8))
set_text(tb2.text_frame, "for Digital Twin Optimization", size=32, color=ACCENT, align=PP_ALIGN.CENTER)
tb3 = sl.shapes.add_textbox(Inches(1.5), Inches(3.25), Inches(10), Inches(0.5))
set_text(tb3.text_frame, "Interim Presentation #2  —  First Runs", size=20, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
div = sl.shapes.add_shape(1, Inches(3.5), Inches(3.9), Inches(6.3), Pt(2))
add_solid_fill(div, ACCENT); div.line.fill.background()
tb4 = sl.shapes.add_textbox(Inches(1), Inches(4.15), Inches(11.2), Inches(0.45))
set_text(tb4.text_frame, "Andreas Demosthenous  &  Marios Olymbios", size=18, color=WHITE, align=PP_ALIGN.CENTER)
tb5 = sl.shapes.add_textbox(Inches(1), Inches(4.65), Inches(11.2), Inches(0.4))
set_text(tb5.text_frame, "EPL 445 – Digital Video Processing  |  University of Cyprus  |  April 23, 2026", size=14, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(sl, 1, TOTAL)


# SLIDE 2 — Agenda
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "  Agenda")
items = [
    "1.  Project Recap & Goals",
    "2.  Dataset & Training Setup",
    "3.  Performance Metrics (Test Set)",
    "4.  Training Curves & Visualizations",
    "5.  Video Inference Pipeline",
    "6.  Next Steps — Final Presentation (25 May 2026)",
]
y = Inches(1.5)
for label in items:
    tb = sl.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(0.52))
    set_text(tb.text_frame, label, size=20, color=UCY_BLUE)
    y += Inches(0.7)
add_footer(sl, 2, TOTAL)


# SLIDE 3 — Recap
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "  Project Recap & Interim 2 Goals",
               subtitle_text="  Topic: Mobile Video Processing & Analysis  (EPL 445, Topic #1)")
goals = [
    ("Scientific Goal", "Intelligent auto-interpretation of traffic video to support city Digital Twin systems via vehicle detection and density estimation."),
    ("Interim 2 Objective", "FIRST RUNS — demonstrate a working, end-to-end pipeline using lightweight models on video frames; verify all metrics on a held-out test set."),
    ("Why MobileViT?", "Light-weight Vision Transformer (ICLR 2022). Balances accuracy and efficiency — designed for mobile/edge deployment."),
    ("Methodological Rigour", "We identified and corrected a sequence-level data-leakage issue, retrained on clean splits — results below are honest."),
]
y = Inches(1.55)
for heading, body in goals:
    box = sl.shapes.add_shape(1, Inches(0.4), y, Inches(12.4), Inches(0.85))
    add_solid_fill(box, LIGHT_GREY); box.line.color.rgb = ACCENT; box.line.width = Pt(1)
    th = sl.shapes.add_textbox(Inches(0.55), y + Pt(4), Inches(2.8), Inches(0.38))
    set_text(th.text_frame, heading + ":", size=14, bold=True, color=UCY_BLUE)
    tb = sl.shapes.add_textbox(Inches(3.35), y + Pt(4), Inches(9.3), Inches(0.55))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, body, size=13, color=DARK_GREY)
    y += Inches(1.0)
add_footer(sl, 3, TOTAL)


# SLIDE 4 — Dataset
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "  Dataset & Training Setup")
left_box = sl.shapes.add_shape(1, Inches(0.3), Inches(1.5), Inches(5.8), Inches(5.2))
add_solid_fill(left_box, WHITE); left_box.line.color.rgb = UCY_BLUE; left_box.line.width = Pt(1.5)
lh = sl.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(5.3), Inches(0.4))
set_text(lh.text_frame, "Dataset: UA-DETRAC", size=16, bold=True, color=UCY_BLUE)
left_items = [
    "Source: Wen et al., 2020 — urban surveillance",
    "4 classes: car | bus | truck | background",
    "100 video sequences total",
    "Sequence-level split (no leakage)",
    "   Train  — 69 sequences,  6,827 patches",
    "   Val    — 15 sequences,    928 patches",
    "   Test   — 16 sequences,  2,614 patches",
    "Crops < 64 px filtered to avoid upscale artefacts",
    "Image size: 256 × 256 input",
    "Augmentation: random flip, rotation, colour jitter",
]
y = Inches(2.1)
for item in left_items:
    tb = sl.shapes.add_textbox(Inches(0.6), y, Inches(5.4), Inches(0.34))
    set_text(tb.text_frame, item, size=12, color=DARK_GREY)
    y += Inches(0.33)
right_box = sl.shapes.add_shape(1, Inches(6.6), Inches(1.5), Inches(6.4), Inches(5.2))
add_solid_fill(right_box, WHITE); right_box.line.color.rgb = ACCENT; right_box.line.width = Pt(1.5)
rh = sl.shapes.add_textbox(Inches(6.85), Inches(1.6), Inches(5.9), Inches(0.4))
set_text(rh.text_frame, "Training Strategy", size=16, bold=True, color=UCY_BLUE)
right_items = [
    ("Staged Fine-Tuning (15 epochs)", True),
    ("Epochs 1–3:  Backbone frozen", False),
    ("         Only classifier head trains (lr = 3e-4)", False),
    ("Epochs 4–15:  Full model unfrozen", False),
    ("         End-to-end fine-tuning (lr = 3e-5)", False),
    ("Optimizer: AdamW  |  Scheduler: CosineAnnealingLR", False),
    ("Loss: Class-weighted Cross-Entropy", False),
    ("Hardware: AMD Radeon 780M (ROCm 7.2)  | bs=4", False),
    ("Best checkpoint chosen by val macro-F1", False),
    ("Best val macro-F1: 94.04%", True),
]
y = Inches(2.1)
for text, bold in right_items:
    tb = sl.shapes.add_textbox(Inches(6.85), y, Inches(6.0), Inches(0.34))
    set_text(tb.text_frame, text, size=12, bold=bold, color=UCY_BLUE if bold else DARK_GREY)
    y += Inches(0.37)
add_footer(sl, 4, TOTAL)


# SLIDE 5 — Metrics
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "  Performance Metrics — Test Set  (n = 2,614)",
               subtitle_text="  Metrics on held-out, sequence-level test set as required by Interim #2 brief")

macros = [
    ("Accuracy",        "96.44%", UCY_BLUE),
    ("Macro F1-Score",  "95.90%", ACCENT),
    ("Macro Precision", "95.69%", RGBColor(0x1A, 0x6A, 0xB0)),
    ("Macro Recall",    "96.15%", RGBColor(0x13, 0x74, 0x8A)),
    ("Macro AUC",       "99.59%", GREEN),
]
card_w = Inches(2.4); card_h = Inches(1.05); gap = Inches(0.18); start_x = Inches(0.25)
for i, (label, value, color) in enumerate(macros):
    metric_card(sl, start_x + i * (card_w + gap), Inches(1.55), card_w, card_h, label, value, color)

th = sl.shapes.add_textbox(Inches(0.35), Inches(2.85), Inches(12.5), Inches(0.38))
set_text(th.text_frame, "Per-Class Results (Test Set)", size=15, bold=True, color=UCY_BLUE)
headers = ["Class", "Precision", "Recall / Sensitivity", "Specificity", "F1-Score", "Support"]
col_widths = [Inches(1.6), Inches(1.8), Inches(2.5), Inches(1.8), Inches(1.8), Inches(1.4)]
col_x = [Inches(0.35)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w + Inches(0.05))
hdr_y = Inches(3.28)
for h, w, x in zip(headers, col_widths, col_x):
    hdr_box = sl.shapes.add_shape(1, x, hdr_y, w, Inches(0.36))
    add_solid_fill(hdr_box, UCY_BLUE); hdr_box.line.fill.background()
    tb = sl.shapes.add_textbox(x + Pt(4), hdr_y + Pt(4), w - Pt(8), Inches(0.3))
    set_text(tb.text_frame, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Car",        "96.04%", "92.92%", "99.20%", "94.45%",  "438"),
    ("Bus",        "99.16%", "97.45%", "99.39%", "98.30%", "1,136"),
    ("Truck",      "93.69%", "96.26%", "97.71%", "94.96%",  "695"),
    ("Background", "93.86%", "97.97%", "99.03%", "95.87%",  "345"),
]
row_colors = [LIGHT_GREY, WHITE, LIGHT_GREY, WHITE]
data_y = hdr_y + Inches(0.38)
for ri, (row_data, bg_c) in enumerate(zip(rows, row_colors)):
    for ci, (val, w, x) in enumerate(zip(row_data, col_widths, col_x)):
        cell = sl.shapes.add_shape(1, x, data_y, w, Inches(0.5))
        add_solid_fill(cell, bg_c)
        cell.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); cell.line.width = Pt(0.5)
        tc = sl.shapes.add_textbox(x + Pt(4), data_y + Pt(6), w - Pt(8), Inches(0.38))
        col = UCY_BLUE if ci == 0 else DARK_GREY
        set_text(tc.text_frame, val, size=13, bold=(ci == 0), color=col, align=PP_ALIGN.CENTER)
    data_y += Inches(0.52)
add_footer(sl, 5, TOTAL)


# SLIDE 6 — Visualizations
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "  Training Curves & Performance Visualizations")
figures = [
    ("training_curves.png",  "Training Curves",   Inches(0.25), Inches(1.45)),
    ("confusion_matrix.png", "Confusion Matrix",  Inches(6.85), Inches(1.45)),
    ("roc_curves.png",       "ROC Curves",        Inches(0.25), Inches(4.15)),
    ("per_class_f1.png",     "Per-Class F1 Score",Inches(6.85), Inches(4.15)),
]
img_w = Inches(6.2); img_h = Inches(2.5)
for fname, caption, lft, tp in figures:
    fpath = os.path.join(FIGURES_DIR, fname)
    cap = sl.shapes.add_textbox(lft, tp - Inches(0.28), img_w, Inches(0.26))
    set_text(cap.text_frame, caption, size=13, bold=True, color=UCY_BLUE)
    if os.path.exists(fpath):
        sl.shapes.add_picture(fpath, lft, tp, img_w, img_h)
add_footer(sl, 6, TOTAL)


# SLIDE 7 — Video pipeline
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "  Video Inference Pipeline")
chal = sl.shapes.add_shape(1, Inches(0.3), Inches(1.55), Inches(12.6), Inches(0.85))
add_solid_fill(chal, RGBColor(0xFF, 0xF3, 0xCD))
chal.line.color.rgb = ORANGE; chal.line.width = Pt(1)
ct = sl.shapes.add_textbox(Inches(0.5), Inches(1.62), Inches(12.2), Inches(0.65))
ct.text_frame.word_wrap = True
set_text(ct.text_frame,
    "Challenge: MobileViT is a patch classifier (trained on cropped objects). "
    "Applied directly to full-resolution frames it has poor localisation.",
    size=13, color=RGBColor(0x80, 0x50, 0x00))
steps = [
    ("1. Frame Extraction", "Video decoded frame-by-frame using OpenCV (with frame_skip)."),
    ("2. Multi-Scale Sliding Window", "Windows of 3 sizes (120, 180, 256 px) stride across each frame; each patch resized to 256×256 for MobileViT."),
    ("3. Non-Maximum Suppression", "Per-class NMS at IoU 0.3, retaining the highest-confidence detection."),
    ("4. Aggregation & Output", "Per-frame class counts → CSV. Aggregated totals → JSON. Annotated bounding boxes → MP4."),
]
y = Inches(2.6)
for title, detail in steps:
    icon = sl.shapes.add_shape(1, Inches(0.35), y + Pt(2), Inches(0.35), Inches(0.35))
    add_solid_fill(icon, ACCENT); icon.line.fill.background()
    tb_h = sl.shapes.add_textbox(Inches(0.85), y, Inches(11.8), Inches(0.32))
    set_text(tb_h.text_frame, title, size=14, bold=True, color=UCY_BLUE)
    tb_d = sl.shapes.add_textbox(Inches(0.85), y + Inches(0.33), Inches(11.8), Inches(0.45))
    tb_d.text_frame.word_wrap = True
    set_text(tb_d.text_frame, detail, size=12, color=DARK_GREY)
    y += Inches(0.97)
res = sl.shapes.add_shape(1, Inches(0.3), Inches(6.45), Inches(12.6), Inches(0.42))
add_solid_fill(res, UCY_BLUE); res.line.fill.background()
rt = sl.shapes.add_textbox(Inches(0.5), Inches(6.48), Inches(12.2), Inches(0.38))
set_text(rt.text_frame,
    "Working video demo with per-frame density estimates  |  Phase 3: replace sliding window with YOLOv8-nano",
    size=14, bold=True, color=WHITE)
add_footer(sl, 7, TOTAL)


# SLIDE 8 — Next Steps
sl = prs.slides.add_slide(blank_layout)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "  Next Steps — Final Presentation (25 May 2026)")
next_steps = [
    ("Detection Upgrade", "Replace sliding-window with YOLOv8-nano for box proposals; pass each crop to MobileViT classifier. Expected >10× speed-up."),
    ("Tracking Engine",   "SORT / ByteTrack to assign persistent IDs across frames; eliminates double-counting for time-series occupancy."),
    ("Web Dashboard",     "FastAPI / Streamlit: upload video, view live counts, ROI lane counting, density heat-maps."),
    ("Real-Time Support", "RTSP stream ingestion for live camera integration; target latency <200 ms/frame on GPU."),
]
y = Inches(1.55)
for title, detail in next_steps:
    box = sl.shapes.add_shape(1, Inches(0.3), y, Inches(12.5), Inches(1.12))
    add_solid_fill(box, WHITE); box.line.color.rgb = ACCENT; box.line.width = Pt(1.2)
    dot = sl.shapes.add_shape(1, Inches(0.35), y + Inches(0.38), Inches(0.12), Inches(0.12))
    add_solid_fill(dot, ACCENT); dot.line.fill.background()
    th = sl.shapes.add_textbox(Inches(0.58), y + Pt(5), Inches(2.5), Inches(0.35))
    set_text(th.text_frame, title + ":", size=14, bold=True, color=UCY_BLUE)
    td = sl.shapes.add_textbox(Inches(3.15), y + Pt(5), Inches(9.5), Inches(0.85))
    td.text_frame.word_wrap = True
    set_text(td.text_frame, detail, size=12, color=DARK_GREY)
    y += Inches(1.22)
ds = sl.shapes.add_shape(1, Inches(0.3), Inches(6.5), Inches(12.6), Inches(0.38))
add_solid_fill(ds, GREEN); ds.line.fill.background()
dt = sl.shapes.add_textbox(Inches(0.5), Inches(6.53), Inches(12.2), Inches(0.33))
set_text(dt.text_frame,
    "Final Project Presentation & Delivery: 25 May 2026, 14:00–17:00  |  Moodle + GitHub Repository",
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_footer(sl, 8, TOTAL)


out_path = os.path.join(os.path.dirname(__file__), "..", "EPL445_Interim2_MobileViT_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {os.path.abspath(out_path)}")
