"""Generate the final-presentation PPTX for EPL445 MobileViT project (25 May 2026)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

UCY_BLUE   = RGBColor(0x00, 0x3A, 0x70)
ACCENT     = RGBColor(0x00, 0x8B, 0xD0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF5, 0xF9)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x1A, 0x8A, 0x42)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)

HERE        = os.path.dirname(__file__)
FIGURES_DIR = os.path.join(HERE, "..", "outputs", "figures")
SCREENS_DIR = os.path.join(HERE, "..")  # for dashboard-*.png if present


# ── Style helpers ────────────────────────────────────────────────
def add_solid_fill(shape, color):
    sp = shape.fill
    sp.solid()
    sp.fore_color.rgb = color


def set_text(tf, text, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, italic=False):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_header_bar(slide, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.35))
    add_solid_fill(bar, UCY_BLUE); bar.line.fill.background()
    tf = bar.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title_text
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = WHITE
    if subtitle_text:
        sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.95), Inches(13), Pt(22))
        set_text(sub.text_frame, "  " + subtitle_text, size=14, color=ACCENT, italic=True)


def add_footer(slide, slide_num, total):
    line = slide.shapes.add_shape(1, Inches(0), Inches(6.95), Inches(13.33), Pt(1.5))
    add_solid_fill(line, ACCENT); line.line.fill.background()
    footer = slide.shapes.add_textbox(Inches(0), Inches(7.0), Inches(13.33), Inches(0.4))
    tf = footer.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (f"EPL 445 – Digital Video Processing  |  University of Cyprus  |  "
                f"Final Presentation, 25 May 2026  |  Slide {slide_num}/{total}")
    run.font.size = Pt(9); run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)


def metric_card(slide, left, top, width, height, label, value, color):
    card = slide.shapes.add_shape(1, left, top, width, height)
    add_solid_fill(card, color); card.line.fill.background()
    tb = slide.shapes.add_textbox(left, top + Pt(6), width, Inches(0.3))
    set_text(tb.text_frame, label, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    tb2 = slide.shapes.add_textbox(left, top + Inches(0.35), width, Inches(0.6))
    set_text(tb2.text_frame, value, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def two_col_card(slide, left, top, width, height, heading, body, accent=ACCENT):
    box = slide.shapes.add_shape(1, left, top, width, height)
    add_solid_fill(box, LIGHT_GREY); box.line.color.rgb = accent; box.line.width = Pt(1)
    th = slide.shapes.add_textbox(left + Pt(8), top + Pt(6), Inches(2.6), Inches(0.4))
    set_text(th.text_frame, heading, size=13, bold=True, color=UCY_BLUE)
    tb = slide.shapes.add_textbox(left + Inches(2.9), top + Pt(6), width - Inches(3.0), height - Pt(12))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, body, size=12, color=DARK_GREY)


# ── Presentation ─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# Slide order
SLIDES = [
    "Title",
    "Agenda",
    "Problem & Motivation",
    "Dataset",
    "MobileViT Architecture",
    "Training Procedure",
    "Detection: Sliding Window vs YOLO",
    "Cross-Frame Tracking (SORT)",
    "Per-Lane Counting + Web Dashboard",
    "Live Dashboard Screenshot",
    "Results — Classification Metrics",
    "Results — Confusion Matrix & ROC",
    "Results — Training & Density",
    "Speed/Accuracy Benchmark",
    "Discussion & Limitations",
    "Future Work",
    "Live Demo",
    "Questions",
]
TOTAL = len(SLIDES)


# ── Slide 1 — Title ──────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, UCY_BLUE); bg.line.fill.background()
stripe = sl.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(0.08))
add_solid_fill(stripe, ACCENT); stripe.line.fill.background()
tb = sl.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.4))
set_text(tb.text_frame, "MobileViT for Real-Time Traffic Monitoring",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.7))
set_text(tb2.text_frame,
         "A Web-Based Classification, Detection, and Tracking System",
         size=22, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
tb3 = sl.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10), Inches(0.4))
set_text(tb3.text_frame, "Final Presentation", size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
div = sl.shapes.add_shape(1, Inches(3.5), Inches(4.1), Inches(6.3), Pt(2))
add_solid_fill(div, ACCENT); div.line.fill.background()
tb4 = sl.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.2), Inches(0.5))
set_text(tb4.text_frame, "Andreas Demosthenous  &  Marios Olympios",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
tb5 = sl.shapes.add_textbox(Inches(1), Inches(4.85), Inches(11.2), Inches(0.4))
set_text(tb5.text_frame,
         "EPL 445 – Digital Video Processing  |  University of Cyprus  |  25 May 2026",
         size=14, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(sl, 1, TOTAL)


# ── Slide 2 — Agenda ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Agenda")
items = [
    "1.  Problem and motivation",
    "2.  Dataset and pre-processing",
    "3.  MobileViT classifier and training",
    "4.  Detection pipeline: sliding window vs YOLO hybrid",
    "5.  Tracking (SORT) and per-lane counting",
    "6.  Web dashboard",
    "7.  Quantitative results",
    "8.  Live demo, discussion, and questions",
]
y = Inches(1.65)
for label in items:
    tb = sl.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(0.52))
    set_text(tb.text_frame, label, size=20, color=UCY_BLUE)
    y += Inches(0.6)
add_footer(sl, 2, TOTAL)


# ── Slide 3 — Problem & Motivation ───────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Problem and Motivation",
               subtitle_text="Live vehicle counts feed digital-twin traffic platforms")
cards = [
    ("Scientific question", "Can a lightweight vision transformer classify traffic patches accurately enough to power a real-time, web-based per-lane counting dashboard on integrated graphics?"),
    ("Why MobileViT", "5.6 M parameters, ImageNet-pretrained, combines convolution with self-attention. Fits in 8 GB shared VRAM where ViT-base does not."),
    ("Why hybrid detection", "MobileViT outputs class probabilities, not boxes. We pair it with YOLOv8-nano for tight proposals and keep the bespoke 4-class decision (car, bus, truck, background) on the classifier."),
    ("End-to-end deliverable", "FastAPI dashboard with video upload, per-lane polygon ROI, density timeline, cross-frame tracking, and live H.264 playback — runs from a single uvicorn command."),
]
y = Inches(1.55)
for heading, body in cards:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(1.15), heading, body)
    y += Inches(1.28)
add_footer(sl, 3, TOTAL)


# ── Slide 4 — Dataset ────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Dataset and Pre-processing",
               subtitle_text="UA-DETRAC, 4 classes, sequence-disjoint splits")
left = sl.shapes.add_shape(1, Inches(0.3), Inches(1.55), Inches(6.4), Inches(5.2))
add_solid_fill(left, WHITE); left.line.color.rgb = ACCENT; left.line.width = Pt(1)
tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(6.0), Inches(5))
tb.text_frame.word_wrap = True
tf = tb.text_frame; tf.clear()
lines = [
    ("Source", "UA-DETRAC via Kaggle mirror (~140 K annotated frames, 100 surveillance sequences)."),
    ("Crops", "Every bounding box larger than 32x32 px, capped at 5 000 per class, ~475 K total."),
    ("Label map", "car -> car, bus -> bus, van -> truck, others -> truck. Negatives sampled from frame corners."),
    ("Pre-processing", "Resize to 256x256, ImageNet normalisation, horizontal flip, +-10 deg rotation, colour jitter."),
    ("Splits", "70/15/15 by SEQUENCE ID (MVI_xxxxx) so no frame leaks across train/val/test."),
    ("Leakage fix", "Initial random splits inflated accuracy to ~99%. Sequence-level splits dropped accuracy to an honest 96%+ — same model, real generalisation."),
]
first = True
for h, b in lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = h + ": "
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = UCY_BLUE
    r2 = p.add_run(); r2.text = b
    r2.font.size = Pt(12); r2.font.color.rgb = DARK_GREY
    p.space_after = Pt(8)

right = sl.shapes.add_shape(1, Inches(7.0), Inches(1.55), Inches(5.95), Inches(5.2))
add_solid_fill(right, WHITE); right.line.color.rgb = ACCENT; right.line.width = Pt(1)
tb = sl.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.6), Inches(0.4))
set_text(tb.text_frame, "Class distribution (after balancing)", size=14, bold=True, color=UCY_BLUE)
counts = [("car", "5 000"), ("bus", "~3 200"), ("truck", "5 000"), ("background", "~2 000")]
y = Inches(2.2)
for cls, n in counts:
    bar = sl.shapes.add_shape(1, Inches(7.2), y, Inches(5.6), Inches(0.55))
    add_solid_fill(bar, ACCENT); bar.line.fill.background()
    tb = sl.shapes.add_textbox(Inches(7.3), y + Pt(6), Inches(2.5), Inches(0.4))
    set_text(tb.text_frame, cls, size=14, bold=True, color=WHITE)
    tb = sl.shapes.add_textbox(Inches(10.5), y + Pt(6), Inches(2.2), Inches(0.4))
    set_text(tb.text_frame, n, size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    y += Inches(0.85)
add_footer(sl, 4, TOTAL)


# ── Slide 5 — MobileViT architecture ────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "MobileViT-S Architecture",
               subtitle_text="Mehta & Rastegari, ICLR 2022 — light-weight, mobile-friendly")
boxes = [
    ("Conv stem", "3x3 conv, downsample to 128x128", ACCENT),
    ("MV2 block x3", "Inverted residuals, depthwise conv, channels x2", UCY_BLUE),
    ("MobileViT block 1", "Local conv -> patch attention -> fold back", GREEN),
    ("MobileViT block 2", "192 channels, 4 transformer layers", GREEN),
    ("MobileViT block 3", "240 channels, 3 transformer layers", GREEN),
    ("Head", "Global avg pool -> linear (4 classes)", ORANGE),
]
y = Inches(1.7)
for i, (h, b, c) in enumerate(boxes):
    card = sl.shapes.add_shape(1, Inches(0.6) + Inches(i * 2.05), y, Inches(1.9), Inches(2.2))
    add_solid_fill(card, c); card.line.fill.background()
    tb = sl.shapes.add_textbox(Inches(0.7) + Inches(i * 2.05), y + Pt(6), Inches(1.7), Inches(0.45))
    set_text(tb.text_frame, h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = sl.shapes.add_textbox(Inches(0.7) + Inches(i * 2.05), y + Inches(0.55), Inches(1.7), Inches(1.5))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, b, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Highlights
hl = sl.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.4), Inches(2.0))
hl.text_frame.word_wrap = True
tf = hl.text_frame; tf.clear()
items = [
    "5.6 M parameters total. Roughly 4 times smaller than ResNet-50 with similar accuracy on UA-DETRAC patches.",
    "Convolutional stages capture local texture (vehicle parts) — transformer blocks model global context within the crop.",
    "Pretrained on ImageNet via timm. Fine-tuned 15 epochs on traffic patches.",
    "Fits in 8 GB shared VRAM on the AMD Radeon 780M iGPU. Trains with batch size 4.",
]
first = True
for it in items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "•  " + it
    r.font.size = Pt(13); r.font.color.rgb = DARK_GREY
    p.space_after = Pt(6)
add_footer(sl, 5, TOTAL)


# ── Slide 6 — Training procedure ─────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Training Procedure",
               subtitle_text="Two-phase fine-tuning, AdamW, cosine schedule")
cards = [
    ("Optimiser", "AdamW with learning rate 3e-4 and weight decay 1e-4. Scheduler: cosine annealing over 15 epochs."),
    ("Two phases", "Phase 1 (epochs 1-3): backbone frozen, classifier head trains. Phase 2 (epochs 4-15): backbone unfrozen, lr reduced by 10x."),
    ("Loss", "Class-weighted cross-entropy (background and bus under-represented after the per-class cap)."),
    ("Hardware", "AMD Radeon 780M iGPU, ROCm 7.2, HSA_OVERRIDE_GFX_VERSION=11.0.0. Batch size 4 due to 8 GB shared VRAM."),
    ("Selection", "Early stopping on macro F1, patience 5. Best val macro F1 = 0.9807 at epoch 13."),
]
y = Inches(1.55)
for h, b in cards:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(0.95), h, b)
    y += Inches(1.05)
add_footer(sl, 6, TOTAL)


# ── Slide 7 — Detection: sliding window vs YOLO ──────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Detection Pipeline",
               subtitle_text="Sliding-window baseline -> YOLOv8-nano + MobileViT hybrid")

# Left card: sliding
slide_card = sl.shapes.add_shape(1, Inches(0.4), Inches(1.55), Inches(6.2), Inches(5.2))
add_solid_fill(slide_card, LIGHT_GREY); slide_card.line.color.rgb = ORANGE; slide_card.line.width = Pt(1.5)
tb = sl.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.8), Inches(0.45))
set_text(tb.text_frame, "Multi-scale sliding window (Interim 2)", size=15, bold=True, color=ORANGE)
items = [
    "Three window sizes: 120, 180, 256 px, 50% stride",
    "Every tile classified by MobileViT",
    "Per-class greedy NMS, IoU = 0.3",
    "~3 fps on 640x360 frames",
    "Overlapping detections, many background tiles",
]
tb = sl.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.6), Inches(4.5))
tb.text_frame.word_wrap = True
tf = tb.text_frame; tf.clear()
for i, it in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + it
    r.font.size = Pt(13); r.font.color.rgb = DARK_GREY
    p.space_after = Pt(8)

# Right card: YOLO hybrid
yolo_card = sl.shapes.add_shape(1, Inches(6.85), Inches(1.55), Inches(6.1), Inches(5.2))
add_solid_fill(yolo_card, LIGHT_GREY); yolo_card.line.color.rgb = GREEN; yolo_card.line.width = Pt(1.5)
tb = sl.shapes.add_textbox(Inches(7.05), Inches(1.7), Inches(5.7), Inches(0.45))
set_text(tb.text_frame, "YOLOv8-nano + MobileViT (Final)", size=15, bold=True, color=GREEN)
items = [
    "YOLO emits ~one tight box per vehicle (COCO 2/3/5/7)",
    "MobileViT classifies each crop in the 4-class taxonomy",
    "Boxes labelled background are discarded",
    "~4 fps steady-state on the same hardware",
    "Single box per object, far fewer duplicates",
    "Same MobileViT backbone reused — no retraining",
]
tb = sl.shapes.add_textbox(Inches(7.05), Inches(2.2), Inches(5.5), Inches(4.5))
tb.text_frame.word_wrap = True
tf = tb.text_frame; tf.clear()
for i, it in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + it
    r.font.size = Pt(13); r.font.color.rgb = DARK_GREY
    p.space_after = Pt(8)
add_footer(sl, 7, TOTAL)


# ── Slide 8 — Tracking ───────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Cross-Frame Tracking (SORT)",
               subtitle_text="Bewley et al., ICIP 2016 — IoU-association tracker")
cards = [
    ("Association", "Per-frame detections matched to existing trackers by IoU. Unmatched detections spawn new trackers; trackers without confirmation for max_age = 5 frames retire."),
    ("Confirmation gating", "min_hits = 2 before a track is reported, suppressing single-frame false positives."),
    ("CSV instrumentation", "Each detection row gets a track_id column. aggregate_counts emits unique_vehicles_by_class alongside raw detection counts."),
    ("Demo impact", "Raw detections over-count: each vehicle appears in N frames. With SORT, the dashboard shows the underlying number of unique vehicles, which is what a digital twin needs."),
]
y = Inches(1.55)
for h, b in cards:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(1.15), h, b)
    y += Inches(1.27)
add_footer(sl, 8, TOTAL)


# ── Slide 9 — Per-lane + dashboard ───────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Per-Lane Counting and Web Dashboard",
               subtitle_text="Polygon ROIs + FastAPI + vanilla HTML/JS, no build step")
cards = [
    ("Polygon ROIs", "Operator draws closed polygons on the first frame. Backend tests detection-centre membership with OpenCV pointPolygonTest. No re-inference required."),
    ("Endpoints", "POST /jobs (upload), GET /jobs/{id}/counts | /timeline | /video | /frame, POST /jobs/{id}/lanes (per-lane), POST /jobs/{id}/roi (rectangle)."),
    ("Frontend", "Single index.html + app.js + style.css. Chart.js via CDN for bar and line charts. HTML5 canvas for ROI and lane editing."),
    ("Robustness", "OpenCV writes mp4v which browsers cannot decode. The /video endpoint lazily re-encodes to H.264 (libx264, yuv420p, +faststart) and the worker pre-warms it after each job."),
]
y = Inches(1.55)
for h, b in cards:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(1.15), h, b)
    y += Inches(1.27)
add_footer(sl, 9, TOTAL)


# ── Slide 10 — Live dashboard screenshot ─────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Live Dashboard",
               subtitle_text="FastAPI + vanilla HTML/JS, no build step")
shot_path = os.path.join(HERE, "..", "docs", "screenshots", "dashboard_full.png")
if os.path.exists(shot_path):
    # Letterbox — original is tall portrait, fit by height
    sl.shapes.add_picture(shot_path, Inches(4.0), Inches(1.55), height=Inches(5.2))
# Annotations on the right side
notes = [
    ("Drag-and-drop upload", "MP4 in, background inference starts, progress bar polls every 3 s"),
    ("Annotated H.264 video", "Lazily re-encoded from OpenCV mp4v so it plays in-browser"),
    ("Counts and timeline", "Bar chart + density line chart via Chart.js (CDN, no build)"),
    ("Rectangle and polygon ROIs", "Filter detections by region or split into per-lane polygons"),
    ("Past jobs table", "Click to re-open any prior run"),
]
y = Inches(1.6)
for h, b in notes:
    tb = sl.shapes.add_textbox(Inches(0.4), y, Inches(3.5), Inches(0.32))
    set_text(tb.text_frame, h, size=13, bold=True, color=UCY_BLUE)
    tb = sl.shapes.add_textbox(Inches(0.4), y + Inches(0.32), Inches(3.5), Inches(0.6))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, b, size=11, color=DARK_GREY)
    y += Inches(1.02)
add_footer(sl, 10, TOTAL)


# ── Slide 11 — Classification metrics ────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Results — Classification Metrics",
               subtitle_text="Sequence-disjoint test set, 2 614 patches, 16 unseen sequences")

# Top: headline metric cards
metric_card(sl, Inches(0.4), Inches(1.6),  Inches(2.4), Inches(1.1), "ACCURACY",   "96.44%", GREEN)
metric_card(sl, Inches(3.0), Inches(1.6),  Inches(2.4), Inches(1.1), "MACRO F1",   "95.90%", UCY_BLUE)
metric_card(sl, Inches(5.6), Inches(1.6),  Inches(2.4), Inches(1.1), "MACRO AUC",  "99.59%", ACCENT)
metric_card(sl, Inches(8.2), Inches(1.6),  Inches(2.4), Inches(1.1), "VAL F1 (best)", "0.9807", ORANGE)

# Per-class table
tb = sl.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(12.4), Inches(0.4))
set_text(tb.text_frame, "Per-class breakdown:", size=14, bold=True, color=UCY_BLUE)
headers = ["Class", "Precision", "Recall", "Specificity", "F1"]
rows = [
    ["car",        "0.960", "0.929", "0.992", "0.944"],
    ["bus",        "0.992", "0.974", "0.994", "0.983"],
    ["truck",      "0.937", "0.963", "0.977", "0.950"],
    ["background", "0.939", "0.980", "0.990", "0.959"],
]
table_top = Inches(3.5)
col_w = [Inches(2.4)] + [Inches(2.5)] * 4
y = table_top
# header
x = Inches(0.4)
for i, h in enumerate(headers):
    cell = sl.shapes.add_shape(1, x, y, col_w[i], Inches(0.45))
    add_solid_fill(cell, UCY_BLUE); cell.line.fill.background()
    tb = sl.shapes.add_textbox(x, y + Pt(6), col_w[i], Inches(0.35))
    set_text(tb.text_frame, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w[i]
y += Inches(0.45)
for r_i, row in enumerate(rows):
    x = Inches(0.4)
    fill = WHITE if r_i % 2 == 0 else LIGHT_GREY
    for i, v in enumerate(row):
        cell = sl.shapes.add_shape(1, x, y, col_w[i], Inches(0.42))
        add_solid_fill(cell, fill); cell.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); cell.line.width = Pt(0.5)
        tb = sl.shapes.add_textbox(x, y + Pt(5), col_w[i], Inches(0.32))
        set_text(tb.text_frame, v, size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
        x += col_w[i]
    y += Inches(0.42)

# Note
tb = sl.shapes.add_textbox(Inches(0.4), Inches(6.0), Inches(12.4), Inches(0.4))
set_text(tb.text_frame,
         "Bus is easiest (large, distinctive). Car is hardest — confused with truck at distance.",
         size=12, color=DARK_GREY, italic=True)
add_footer(sl, 11, TOTAL)


# ── Slide 12 — Figures: confusion + ROC ──────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Results — Confusion Matrix and ROC")
cm = os.path.join(FIGURES_DIR, "confusion_matrix.png")
roc = os.path.join(FIGURES_DIR, "roc_curves.png")
if os.path.exists(cm):
    sl.shapes.add_picture(cm,  Inches(0.4), Inches(1.6), width=Inches(6.2))
if os.path.exists(roc):
    sl.shapes.add_picture(roc, Inches(6.8), Inches(1.6), width=Inches(6.2))
tb = sl.shapes.add_textbox(Inches(0.4), Inches(6.4), Inches(12.4), Inches(0.4))
set_text(tb.text_frame,
         "Mistakes concentrate on the car ↔ truck boundary. All four ROC curves crowd the top-left corner.",
         size=13, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)
add_footer(sl, 12, TOTAL)


# ── Slide 13 — Figures: training + density ───────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Results — Training Curves and Density Over Time")
tr = os.path.join(FIGURES_DIR, "training_curves.png")
de = os.path.join(FIGURES_DIR, "density_plot.png")
if os.path.exists(tr):
    sl.shapes.add_picture(tr, Inches(0.4), Inches(1.6), width=Inches(6.2))
if os.path.exists(de):
    sl.shapes.add_picture(de, Inches(6.8), Inches(1.6), width=Inches(6.2))
tb = sl.shapes.add_textbox(Inches(0.4), Inches(6.4), Inches(12.4), Inches(0.4))
set_text(tb.text_frame,
         "Backbone unfreezes at epoch 4 — clear inflection. Density peaks correspond to platoon arrivals.",
         size=13, color=DARK_GREY, italic=True, align=PP_ALIGN.CENTER)
add_footer(sl, 13, TOTAL)


# ── Slide 14 — Speed/accuracy benchmark ─────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Speed / Accuracy Benchmark",
               subtitle_text="MobileViT classifier on AMD Radeon 780M (iGPU)")
headers = ["Input size", "Latency (ms)", "Accuracy", "Use case"]
rows = [
    ["128 x 128", "~5  ms",  "94.5%",  "Edge / battery-constrained"],
    ["192 x 192", "~7  ms",  "95.8%",  "Real-time stream"],
    ["256 x 256", "~13 ms",  "96.44%", "Production (default)"],
]
y = Inches(1.8)
x = Inches(1.5); col_w = [Inches(2.2), Inches(2.2), Inches(2.2), Inches(3.7)]
for i, h in enumerate(headers):
    cell = sl.shapes.add_shape(1, x, y, col_w[i], Inches(0.6))
    add_solid_fill(cell, UCY_BLUE); cell.line.fill.background()
    tb = sl.shapes.add_textbox(x, y + Pt(10), col_w[i], Inches(0.4))
    set_text(tb.text_frame, h, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w[i]
y += Inches(0.6)
for r_i, row in enumerate(rows):
    x = Inches(1.5)
    fill = WHITE if r_i % 2 == 0 else LIGHT_GREY
    for i, v in enumerate(row):
        cell = sl.shapes.add_shape(1, x, y, col_w[i], Inches(0.6))
        add_solid_fill(cell, fill); cell.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); cell.line.width = Pt(0.5)
        tb = sl.shapes.add_textbox(x, y + Pt(10), col_w[i], Inches(0.4))
        set_text(tb.text_frame, v, size=14, color=DARK_GREY, align=PP_ALIGN.CENTER)
        x += col_w[i]
    y += Inches(0.6)
# Throughput note
tb = sl.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.4))
tb.text_frame.word_wrap = True
set_text(tb.text_frame,
         "End-to-end pipeline (YOLO + MobileViT + SORT + annotation): "
         "~4 fps on 4K frames downscaled to 640x360 for inference, vs ~3 fps for the sliding-window baseline. "
         "First-frame YOLO compile adds a one-off ~15 s warmup.",
         size=15, color=DARK_GREY, italic=True)
add_footer(sl, 14, TOTAL)


# ── Slide 15 — Discussion & Limitations ─────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, LIGHT_GREY); bg.line.fill.background()
add_header_bar(sl, "Discussion and Limitations")
cards = [
    ("What worked", "Decoupling detection from classification: YOLO emits high-quality boxes from COCO pretraining; MobileViT handles the bespoke 4-class decision with the background filter."),
    ("Why MobileViT, not ResNet-50", "Roughly 4x smaller, comparable accuracy on UA-DETRAC, fits in 8 GB shared VRAM. Self-attention helps with partial occlusion."),
    ("Limitation: no temporal model", "SORT only does IoU association. Full mutual occlusion collapses two tracks into one — a problem for dense intersections."),
    ("Limitation: motorcycles", "Our 4 classes do not include motorcycles. YOLO class-3 boxes get pushed through MobileViT, which usually labels them car or background."),
    ("Limitation: iGPU batch size", "AMD Radeon 780M with 8 GB shared VRAM forces batch size 4 in training. A discrete GPU would lift this and speed up training 3-4 x."),
]
y = Inches(1.55)
for h, b in cards:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(0.95), h, b)
    y += Inches(1.04)
add_footer(sl, 15, TOTAL)


# ── Slide 16 — Future work ──────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, WHITE); bg.line.fill.background()
add_header_bar(sl, "Future Work")
items = [
    ("ByteTrack", "Stronger association for longer occlusions, with re-identification embedding."),
    ("Motion-guided sliding", "Restrict the fallback sliding window to motion masks from frame differencing."),
    ("End-to-end YOLO comparison", "Fine-tune YOLOv8 on UA-DETRAC for an apples-to-apples speed-accuracy comparison against the hybrid."),
    ("INT8 quantisation", "ONNX export + INT8 quantisation for Jetson Nano-class edge deployment."),
    ("Temporal voting", "Aggregate MobileViT predictions across short sliding windows of tracker frames for stable per-vehicle labels."),
]
y = Inches(1.6)
for h, b in items:
    two_col_card(sl, Inches(0.4), y, Inches(12.4), Inches(0.85), h, b)
    y += Inches(0.95)
add_footer(sl, 16, TOTAL)


# ── Slide 17 — Live demo ────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, UCY_BLUE); bg.line.fill.background()
stripe = sl.shapes.add_shape(1, Inches(0), Inches(2.7), Inches(13.33), Inches(0.08))
add_solid_fill(stripe, ACCENT); stripe.line.fill.background()
tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.4))
set_text(tb.text_frame, "Live Demo", size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb = sl.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.6))
set_text(tb.text_frame, "uvicorn src.app.main:app  →  http://localhost:8000",
         size=22, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
script = [
    "1.  Upload sample_traffic.mp4 (or stream from webcam)",
    "2.  Watch live progress, then the annotated H.264 video",
    "3.  Inspect per-class bar chart + density timeline",
    "4.  Draw two polygon lanes on the first frame — show per-lane counts",
    "5.  Toggle to sliding-window detector for the speed comparison",
]
y = Inches(3.9)
for s in script:
    tb = sl.shapes.add_textbox(Inches(2.0), y, Inches(9.5), Inches(0.45))
    set_text(tb.text_frame, s, size=16, color=WHITE)
    y += Inches(0.55)
add_footer(sl, 17, TOTAL)


# ── Slide 18 — Q&A ──────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
add_solid_fill(bg, UCY_BLUE); bg.line.fill.background()
stripe = sl.shapes.add_shape(1, Inches(0), Inches(3.0), Inches(13.33), Inches(0.08))
add_solid_fill(stripe, ACCENT); stripe.line.fill.background()
tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.4))
set_text(tb.text_frame, "Thank you — questions?", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb = sl.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6))
set_text(tb.text_frame, "Andreas Demosthenous & Marios Olympios",
         size=22, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
tb = sl.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5))
set_text(tb.text_frame, "{ademos05, molymp01}@ucy.ac.cy", size=18, color=WHITE, align=PP_ALIGN.CENTER)
tb = sl.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5))
set_text(tb.text_frame, "Repository: github.com/AndrewDemsDS/EPL445_MobileViT",
         size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_footer(sl, 18, TOTAL)


# ── Save ────────────────────────────────────────────────────────
out = os.path.join(HERE, "..", "EPL445_Final_MobileViT_Presentation.pptx")
prs.save(out)
print(f"Wrote {out} ({os.path.getsize(out):,} bytes, {len(prs.slides)} slides)")
